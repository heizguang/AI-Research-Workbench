"""
PPT 服务 V2 - 基于 python-pptx 的原生可编辑 PPTX 生成

所有幻灯片使用原生形状（文本框、色块、装饰线），
每个元素都可以在 PowerPoint/WPS 中直接编辑。

流程：报告内容 → 提取大纲 → Agent理解生成PPT内容 → 生成PPTX
"""

import os
import json
import time
import logging
from typing import Dict, Any, List, Optional
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)


# 导入LLM模块
def _get_llm():
    """延迟导入LLM模块并创建实例"""
    import openai
    from dotenv import load_dotenv

    # 获取 backend 目录并加载 .env
    backend_dir = Path(__file__).parent.parent.parent
    env_path = backend_dir / ".env"
    if env_path.exists():
        load_dotenv(str(env_path))

    # 从环境变量读取配置
    smart_llm = os.getenv("SMART_LLM", "openai:mimo-v2.5-pro")
    if ":" in smart_llm:
        model_name = smart_llm.split(":")[1]
    else:
        model_name = smart_llm

    # 创建 LLM 配置类
    class LLMConfig:
        def __init__(self, model=None, temperature=0.7, max_tokens=4000):
            self.model = model or model_name
            self.temperature = temperature
            self.max_tokens = max_tokens
            self.top_p = 1.0
            self.frequency_penalty = 0.0
            self.presence_penalty = 0.0

    # 创建 LLM 实例
    class LLMInstance:
        def __init__(self):
            self.config = LLMConfig()
            self.client = openai.AsyncOpenAI(
                api_key=os.getenv("OPENAI_API_KEY"),
                base_url=os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1")
            )

        async def generate_text(self, prompt, system_prompt=None, config=None):
            cfg = config or self.config
            messages = []
            if system_prompt:
                messages.append({"role": "system", "content": system_prompt})
            messages.append({"role": "user", "content": prompt})

            response = await self.client.chat.completions.create(
                model=cfg.model,
                messages=messages,
                temperature=cfg.temperature,
                max_tokens=cfg.max_tokens
            )
            return response.choices[0].message.content

    return LLMInstance(), LLMConfig


async def _agent_generate_ppt_content(outline: List[Dict[str, Any]], report_content: str) -> List[Dict[str, Any]]:
    """使用 Agent 理解大纲并生成适合 PPT 展示的内容

    Args:
        outline: 提取的大纲结构
        report_content: 原始报告内容（供 Agent 参考）

    Returns:
        Agent 生成的 PPT 内容结构
    """
    llm, LLMConfig = _get_llm()


# ───────────────────── 风格配置 ─────────────────────

STYLES = {
    "professional": {
        "primary":   RGBColor(0, 82, 155),
        "secondary": RGBColor(0, 120, 215),
        "accent":    RGBColor(255, 152, 0),
        "dark":      RGBColor(33, 33, 33),
        "light":     RGBColor(255, 255, 255),
        "gray":      RGBColor(117, 117, 117),
        "bg":        RGBColor(255, 255, 255),
        "light_bg":  RGBColor(245, 247, 250),
    },
    "creative": {
        "primary":   RGBColor(255, 87, 34),
        "secondary": RGBColor(255, 152, 0),
        "accent":    RGBColor(76, 175, 80),
        "dark":      RGBColor(33, 33, 33),
        "light":     RGBColor(255, 255, 255),
        "gray":      RGBColor(117, 117, 117),
        "bg":        RGBColor(255, 255, 255),
        "light_bg":  RGBColor(255, 249, 245),
    },
    "minimal": {
        "primary":   RGBColor(33, 33, 33),
        "secondary": RGBColor(66, 66, 66),
        "accent":    RGBColor(200, 200, 200),
        "dark":      RGBColor(33, 33, 33),
        "light":     RGBColor(255, 255, 255),
        "gray":      RGBColor(150, 150, 150),
        "bg":        RGBColor(255, 255, 255),
        "light_bg":  RGBColor(248, 248, 248),
    },
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """将 #RRGGBB 转为 RGBColor"""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ───────────────────── 工具函数 ─────────────────────

def _add_shape(slide, shape_type, left, top, width, height, color):
    """添加装饰形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=Pt(16),
                 color=None, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """添加文本框"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color or RGBColor(33, 33, 33)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return textbox


def _add_bullet_list(slide, left, top, width, height, items: List[str],
                     font_size=Pt(16), color=None, line_spacing=Pt(28),
                     font_name="Microsoft YaHei"):
    """添加要点列表"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  •  {item}"
        p.font.size = font_size
        p.font.color.rgb = color or RGBColor(33, 33, 33)
        p.font.name = font_name
        p.space_after = Pt(12)
        p.space_before = Pt(4)
        p.line_spacing = line_spacing

    return textbox


def _get_slide_layout(prs):
    """获取空白布局"""
    return prs.slide_layouts[6]


# ───────────────────── 布局生成器 ─────────────────────

def _cover_slide(prs, title: str, subtitle: str, c) -> None:
    """封面页 - 色块背景 + 大标题"""
    slide = prs.slides.add_slide(_get_slide_layout(prs))
    W, H = prs.slide_width, prs.slide_height

    # 底部色块
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), W, Inches(3.2), c["primary"])

    # 装饰线
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(1.5), Inches(3.35), Inches(2), Pt(4), c["accent"])

    # 主标题
    _add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(1.5),
                 title, Pt(44), c["light"], bold=True, alignment=PP_ALIGN.LEFT)

    # 副标题
    if subtitle:
        _add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(0.8),
                     subtitle, Pt(20), RGBColor(200, 220, 240), alignment=PP_ALIGN.LEFT)

    # 底部信息
    _add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
                 "多智能体报告生成系统 · AI Generated", Pt(14), c["gray"],
                 alignment=PP_ALIGN.LEFT)


def _toc_slide(prs, title: str, items: List[str], c) -> None:
    """目录页 - 左侧色条 + 编号列表"""
    slide = prs.slides.add_slide(_get_slide_layout(prs))
    W = prs.slide_width

    # 顶部色条
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), W, Inches(0.08), c["primary"])

    # 左侧装饰
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0.08), Inches(0.4), Inches(7.42), c["secondary"])

    # 标题
    _add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.9),
                 title, Pt(28), c["primary"], bold=True)

    # 下划线
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(1), Inches(1.45), Inches(1.5), Pt(3), c["accent"])

    # 编号列表
    for i, item in enumerate(items):
        y = 2.0 + i * 0.55
        # 编号圆形
        circle = _add_shape(slide, MSO_SHAPE.OVAL,
                            Inches(1.2), Inches(y), Inches(0.35), Inches(0.35),
                            c["primary"])
        circle.text_frame.paragraphs[0].text = str(i + 1)
        circle.text_frame.paragraphs[0].font.size = Pt(12)
        circle.text_frame.paragraphs[0].font.color.rgb = c["light"]
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.paragraphs[0].font.name = "Microsoft YaHei"

        # 文字
        _add_textbox(slide, Inches(1.8), Inches(y), Inches(9), Inches(0.35),
                     item, Pt(16), c["dark"])


def _content_slide(prs, title: str, items: List[str], c,
                   index: int, total: int) -> None:
    """内容页 - 标题 + 要点列表 + 页码"""
    slide = prs.slides.add_slide(_get_slide_layout(prs))
    W = prs.slide_width

    # 顶部色条
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), W, Inches(0.08), c["primary"])

    # 左侧装饰
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0.08), Inches(0.4), Inches(7.42), c["secondary"])

    # 页码
    _add_textbox(slide, Inches(12), Inches(6.8), Inches(1), Inches(0.5),
                 f"{index}/{total - 2}", Pt(11), c["gray"],
                 alignment=PP_ALIGN.RIGHT)

    # 标题
    _add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.9),
                 title, Pt(28), c["primary"], bold=True)

    # 下划线
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(1), Inches(1.45), Inches(1.5), Pt(3), c["accent"])

    # 要点列表
    _add_bullet_list(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(5),
                     items, Pt(16), c["dark"])


def _two_column_slide(prs, title: str, left_title: str, left_items: List[str],
                      right_title: str, right_items: List[str], c,
                      index: int, total: int) -> None:
    """双栏页 - 左右两栏对比"""
    slide = prs.slides.add_slide(_get_slide_layout(prs))
    W = prs.slide_width

    # 顶部色条
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), W, Inches(0.08), c["primary"])

    # 左侧装饰
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0.08), Inches(0.4), Inches(7.42), c["secondary"])

    # 页码
    _add_textbox(slide, Inches(12), Inches(6.8), Inches(1), Inches(0.5),
                 f"{index}/{total - 2}", Pt(11), c["gray"],
                 alignment=PP_ALIGN.RIGHT)

    # 标题
    _add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.9),
                 title, Pt(28), c["primary"], bold=True)

    # 下划线
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(1), Inches(1.45), Inches(1.5), Pt(3), c["accent"])

    # 左栏标题
    _add_textbox(slide, Inches(1), Inches(1.8), Inches(5), Inches(0.5),
                 left_title, Pt(18), c["secondary"], bold=True)

    # 左栏内容
    _add_bullet_list(slide, Inches(1), Inches(2.3), Inches(5), Inches(4),
                     left_items, Pt(14), c["dark"])

    # 分隔线
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(6.2), Inches(1.8), Pt(2), Inches(4.5), c["accent"])

    # 右栏标题
    _add_textbox(slide, Inches(6.5), Inches(1.8), Inches(5), Inches(0.5),
                 right_title, Pt(18), c["secondary"], bold=True)

    # 右栏内容
    _add_bullet_list(slide, Inches(6.5), Inches(2.3), Inches(5), Inches(4),
                     right_items, Pt(14), c["dark"])


def _ending_slide(prs, title: str, subtitle: str, c) -> None:
    """结束页 - 全屏色块 + 感谢"""
    slide = prs.slides.add_slide(_get_slide_layout(prs))
    W, H = prs.slide_width, prs.slide_height

    # 全屏背景
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), W, H, c["primary"])

    # 装饰线
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(5.5), Inches(2.8), Inches(2.3), Pt(3), c["accent"])

    # 主标题
    _add_textbox(slide, Inches(1), Inches(2), Inches(11.333), Inches(1.2),
                 title, Pt(48), c["light"], bold=True, alignment=PP_ALIGN.CENTER)

    # 副内容
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(3.3), Inches(11.333), Inches(1.5),
                     subtitle, Pt(18), RGBColor(200, 220, 240),
                     alignment=PP_ALIGN.CENTER)

    # 底部
    _add_textbox(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
                 "多智能体报告生成系统 · AI Generated", Pt(12),
                 RGBColor(150, 170, 190), alignment=PP_ALIGN.CENTER)


# ───────────────────── 大纲提取 ─────────────────────

def _extract_outline(report_content: str) -> List[Dict[str, Any]]:
    """从报告中提取 PPT 大纲

    支持 Markdown 结构：
    - # 标题 → 报告标题（用于封面）
    - ## 章节 → 一级章节
    - ### 子章节 → 二级章节（单独生成幻灯片）
    - 要点（- 或 * 或数字列表）→ 内容要点

    返回 [{"title": str, "content": [str], "type": "cover"|"toc"|"content"|"ending"}]
    """
    lines = report_content.split("\n")

    # 1. 提取报告标题（第一个 # 标题）
    report_title = ""
    for line in lines:
        if line.startswith("# ") and not line.startswith("## "):
            report_title = line[2:].strip()
            break

    # 2. 按 ## 和 ### 拆分章节
    sections = []
    current_section = None
    current_subsection = None

    for line in lines:
        stripped = line.strip()

        if stripped.startswith("## ") and not stripped.startswith("### "):
            # 新的 ## 章节
            if current_section:
                sections.append(current_section)
            title = stripped[3:].strip()
            # 跳过"参考来源"章节
            if "参考来源" in title or "参考文献" in title:
                current_section = {"title": title, "content": [], "skip": True}
                continue
            current_section = {"title": title, "content": [], "subsections": []}
            current_subsection = None

        elif stripped.startswith("### "):
            # 新的 ### 子章节
            if current_section and not current_section.get("skip"):
                title = stripped[4:].strip()
                current_subsection = {"title": title, "content": []}
                current_section["subsections"].append(current_subsection)

        elif current_section and not current_section.get("skip"):
            # 普通内容行
            if stripped and not stripped.startswith("## "):
                # 提取要点（- 或 * 或数字列表）
                if stripped.startswith(("- ", "* ")):
                    point = stripped[2:].strip()
                    if current_subsection:
                        current_subsection["content"].append(point)
                    else:
                        current_section["content"].append(point)
                elif stripped[0].isdigit() and ". " in stripped[:5]:
                    point = stripped.split(". ", 1)[-1].strip() if ". " in stripped else stripped
                    if current_subsection:
                        current_subsection["content"].append(point)
                    else:
                        current_section["content"].append(point)
                elif not stripped.startswith("#") and not stripped.startswith("|"):
                    # 非标题、非表格的普通段落
                    if len(stripped) > 10:
                        point = stripped
                        if current_subsection:
                            current_subsection["content"].append(point)
                        else:
                            current_section["content"].append(point)

    if current_section:
        sections.append(current_section)

    if not sections:
        return [{"title": report_title or "报告", "content": [], "type": "cover"},
                {"title": "感谢", "content": "感谢观看", "type": "ending"}]

    # 3. 构建大纲
    result = []

    # 封面
    first_title = sections[0]["title"] if sections else "报告"
    result.append({
        "title": report_title or first_title,
        "content": [sections[0]["title"]] if len(sections) > 1 else [],
        "type": "cover"
    })

    # 目录（所有章节标题）
    toc_items = [s["title"] for s in sections]
    result.append({
        "title": "目录",
        "content": toc_items,
        "type": "toc"
    })

    # 内容页
    for section in sections:
        # 跳过"参考来源"等章节
        if "参考" in section["title"]:
            continue

        # 如果有子章节，每个子章节单独一页
        if section.get("subsections"):
            for sub in section["subsections"]:
                content = sub["content"][:6] if sub["content"] else ["（详见报告）"]
                result.append({
                    "title": f"{section['title']} - {sub['title']}",
                    "content": content,
                    "type": "content"
                })
            # 如果章节本身也有内容，再加一页
            if section["content"]:
                result.append({
                    "title": section["title"],
                    "content": section["content"][:6],
                    "type": "content"
                })
        else:
            # 无子章节，直接用章节内容
            content = section["content"][:6] if section["content"] else ["（详见报告）"]
            result.append({
                "title": section["title"],
                "content": content,
                "type": "content"
            })

    # 结束页
    result.append({
        "title": "感谢",
        "content": "感谢观看，欢迎提问",
        "type": "ending"
    })

    return result[:20]


# ───────────────────── Agent 内容生成 ─────────────────────

async def _agent_generate_ppt_content(outline: List[Dict[str, Any]], report_content: str) -> List[Dict[str, Any]]:
    """使用 Agent 理解大纲并生成适合 PPT 展示的内容

    Args:
        outline: 提取的大纲结构
        report_content: 原始报告内容（供 Agent 参考）

    Returns:
        Agent 生成的 PPT 内容结构
    """
    llm, LLMConfig = _get_llm()

    # 构建提示词
    system_prompt = """你是一个专业的 PPT 内容策划专家。你的任务是根据提供的报告大纲，生成适合在 PPT 中展示的精炼内容。

要求：
1. 每页 PPT 的要点控制在 3-5 个，每个要点简洁有力（15-30字）
2. 内容要高度概括，突出关键信息和数据
3. 使用专业的商务语言，适合演示场景
4. 封面页需要一个吸引人的副标题
5. 目录页只列出章节标题
6. 结束页要有感谢语和联系方式提示

输出格式（JSON）：
```json
{
  "pages": [
    {
      "type": "cover",
      "title": "报告主标题",
      "subtitle": "一句话概括报告核心价值"
    },
    {
      "type": "toc",
      "title": "目录",
      "items": ["章节1标题", "章节2标题", ...]
    },
    {
      "type": "content",
      "title": "章节标题",
      "points": ["要点1", "要点2", "要点3"]
    },
    {
      "type": "ending",
      "title": "感谢",
      "subtitle": "感谢观看，期待交流"
    }
  ]
}
```

注意：
- 只输出 JSON，不要有其他内容
- content 类型的页面保持 3-5 个要点
- 如果原大纲内容较多，进行精炼和概括
- 确保内容逻辑连贯，适合演讲展示"""

    # 构建用户提示
    outline_desc = []
    for item in outline:
        if item["type"] == "cover":
            outline_desc.append(f"[封面] {item['title']}")
        elif item["type"] == "toc":
            outline_desc.append(f"[目录] {', '.join(item.get('content', []))}")
        elif item["type"] == "content":
            points = item.get('content', [])
            outline_desc.append(f"[内容] {item['title']} - 要点: {', '.join(points[:3])}...")
        elif item["type"] == "ending":
            outline_desc.append(f"[结束] {item['title']}")

    user_prompt = f"""请根据以下报告大纲，生成适合 PPT 展示的精炼内容：

=== 报告大纲 ===
{chr(10).join(outline_desc)}

=== 原始报告内容（供参考） ===
{report_content}

请生成 PPT 内容（JSON 格式）："""

    try:
        logger.info(f"[PPT Agent] 开始调用 Agent 生成 PPT 内容...")
        start_time = time.time()

        # 调用 LLM
        response = await llm.generate_text(
            prompt=user_prompt,
            system_prompt=system_prompt,
            config=LLMConfig(temperature=0.7, max_tokens=4000)
        )

        elapsed = time.time() - start_time
        logger.info(f"[PPT Agent] Agent 响应完成 ({elapsed:.2f}s)")
        logger.info(f"[PPT Agent] 响应长度: {len(response)} 字符")

        # 解析 JSON
        # 提取 JSON 部分（处理可能的 markdown 代码块）
        json_str = response
        if "```json" in json_str:
            json_str = json_str.split("```json")[1].split("```")[0]
        elif "```" in json_str:
            json_str = json_str.split("```")[1].split("```")[0]

        result = json.loads(json_str.strip())
        pages = result.get("pages", [])

        if not pages:
            logger.warning("[PPT Agent] Agent 返回空内容，使用原始大纲")
            return outline

        logger.info(f"[PPT Agent] 成功解析 {len(pages)} 页内容")
        return pages

    except json.JSONDecodeError as e:
        logger.error(f"[PPT Agent] JSON 解析失败: {e}")
        logger.error(f"[PPT Agent] 原始响应: {response[:500]}")
        return outline
    except Exception as e:
        logger.error(f"[PPT Agent] Agent 调用失败: {e}")
        return outline


# ───────────────────── 主服务 ─────────────────────

class PPTServiceV2:
    """PPT 服务 V2 - 原生可编辑 PPTX 生成"""

    def __init__(self, output_dir: str = "./data/ppt"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    async def create_ppt_from_report(
        self,
        report_content: str,
        template: str = "default",
        style: str = "professional",
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """从报告内容创建可编辑 PPT"""
        options = options or {}
        start_time = time.time()

        logger.info(f"[PPT V2] ========== 开始生成 PPT ==========")
        logger.info(f"[PPT V2] 配置 | 模板: {template} | 风格: {style}")
        logger.info(f"[PPT V2] 输入 | 报告长度: {len(report_content)} 字符")

        # 1. 提取大纲
        step_start = time.time()
        outline = _extract_outline(report_content)
        logger.info(f"[PPT V2] 大纲提取完成 ({time.time()-step_start:.2f}s): {len(outline)} 页")
        for i, item in enumerate(outline):
            content_count = len(item.get('content', [])) if isinstance(item.get('content'), list) else 0
            logger.info(f"[PPT V2]   第{i+1}页 [{item['type']}]: {item['title']} | 要点: {content_count} 个")

        # 2. Agent 理解大纲并生成 PPT 内容
        step_start = time.time()
        logger.info(f"[PPT V2] 开始 Agent 内容生成...")
        ppt_content = await _agent_generate_ppt_content(outline, report_content)
        logger.info(f"[PPT V2] Agent 内容生成完成 ({time.time()-step_start:.2f}s): {len(ppt_content)} 页")

        # 3. 选择风格配置
        c = STYLES.get(style, STYLES["professional"])

        # 处理自定义颜色
        if options.get("color_scheme") == "custom" and options.get("custom_colors"):
            custom = options["custom_colors"]
            if custom.get("primary"):
                c["primary"] = _hex_to_rgb(custom["primary"])
            if custom.get("secondary"):
                c["secondary"] = _hex_to_rgb(custom["secondary"])
            if custom.get("accent"):
                c["accent"] = _hex_to_rgb(custom["accent"])

        # 4. 创建 PPTX
        step_start = time.time()
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        total = len(ppt_content)
        for i, item in enumerate(ppt_content):
            page_start = time.time()
            page_type = item["type"]
            title = item["title"]
            # Agent 生成的内容可能是 points 或 items 或 content
            content = item.get("points", item.get("items", item.get("content", [])))

            logger.info(f"[PPT V2] 生成第 {i+1}/{total} 页 [{page_type}]: {title}")

            if page_type == "cover":
                subtitle = item.get("subtitle", "")
                if not subtitle and isinstance(content, list) and content:
                    subtitle = content[0]
                _cover_slide(prs, title, subtitle, c)

            elif page_type == "toc":
                _toc_slide(prs, title, content if isinstance(content, list) else [], c)

            elif page_type == "content":
                _content_slide(prs, title, content if isinstance(content, list) else [],
                               c, i, total)

            elif page_type == "ending":
                subtitle = item.get("subtitle", "")
                if not subtitle:
                    subtitle = content if isinstance(content, str) else (
                        content[0] if isinstance(content, list) and content else "感谢观看"
                    )
                _ending_slide(prs, title, subtitle, c)

            elapsed = time.time() - page_start
            logger.info(f"[PPT V2]   第{i+1}页完成 ({elapsed:.2f}s)")

        # 5. 保存
        date_str = datetime.now().strftime('%Y%m%d_%H%M%S')
        pptx_path = self.output_dir / f"presentation_{template}_{style}_{date_str}.pptx"
        prs.save(str(pptx_path))

        pptx_size = os.path.getsize(str(pptx_path))
        total_time = time.time() - start_time
        logger.info(f"[PPT V2] 步骤完成 | 文件: {pptx_path}")
        logger.info(f"[PPT V2] 输出 | 大小: {pptx_size/1024:.1f} KB | 总耗时: {total_time:.2f}s")
        logger.info(f"[PPT V2] ========== PPT 生成结束 ==========")

        return str(pptx_path)

    async def get_ppt_template(self, template_name: str) -> Dict[str, Any]:
        templates = {
            "default": {"name": "默认模板", "description": "简洁专业的默认模板", "style": "professional"},
            "business": {"name": "商务模板", "description": "适合商务汇报的模板", "style": "professional"},
            "creative": {"name": "创意模板", "description": "富有创意的设计模板", "style": "creative"},
            "minimal": {"name": "简约模板", "description": "简约风格的模板", "style": "minimal"},
        }
        return templates.get(template_name, templates["default"])

    async def list_templates(self) -> List[Dict[str, Any]]:
        return [
            await self.get_ppt_template("default"),
            await self.get_ppt_template("business"),
            await self.get_ppt_template("creative"),
            await self.get_ppt_template("minimal"),
        ]
