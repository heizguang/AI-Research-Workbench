"""
PPT Skill Service - 基于 ppt-polished-deck-collab skill 脚本链路的 PPT 生成服务

流程：报告内容 → 初始化 workspace → LLM 生成 narrative → 派生 slide_specs → python-pptx 构建 → 质量检查 → 导出 PPTX

与 PPTMasterService 的区别：
- 使用 Skill 的标准脚本链路（init_deck_workspace, derive_slide_specs 等）
- 用 python-pptx 原生构建，输出完全可编辑的 PPTX
- 包含 Skill 定义的质量检查（package_preflight, structure_precheck）
"""

import os
import json
import re
import time
import logging
import subprocess
import shutil
from typing import Dict, Any, List, Optional
from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)

# 路径配置
PROJECT_ROOT = Path(__file__).parent.parent.parent.parent
SKILL_DIR = PROJECT_ROOT / "skills" / "ppt-polished-deck-collab"
SCRIPTS_DIR = SKILL_DIR / "scripts"
REFERENCES_DIR = SKILL_DIR / "references"

# 输出目录
OUTPUT_DIR = PROJECT_ROOT / "data" / "ppt"


def _get_llm():
    """获取 LLM 客户端和模型名 — 复用主 LLM 模块，保持 provider 配置一致"""
    from core.llm import get_llm as _get_main_llm
    llm = _get_main_llm()
    # 使用主 LLM 的 smart 模型配置（已通过 LLM_PROVIDERS 解析）
    model_name = llm.config.model
    client = llm.providers[0].client
    return client, model_name


def _set_font(font, font_name: str, font_size: float = None, bold: bool = None, color: tuple = None):
    """设置字体，同时设置 Latin 和 East Asian 字体"""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from lxml import etree

    font.name = font_name

    # 设置 East Asian 字体 - 通过 run 的 XML 元素
    try:
        r_element = font._element
        rPr = r_element.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.SubElement(r_element, qn('a:rPr'))
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font_name)
    except Exception:
        pass

    if font_size is not None:
        font.size = Pt(font_size)
    if bold is not None:
        font.bold = bold
    if color is not None:
        font.color.rgb = RGBColor(*color)


def _adaptive_font_size(text: str, base_size: float, box_width_in: float, box_height_in: float) -> float:
    """根据文字长度和容器大小自适应计算字号

    Args:
        text: 文本内容
        base_size: 基准字号
        box_width_in: 容器宽度（英寸）
        box_height_in: 容器高度（英寸）
    """
    text_len = len(text)

    # 估算每行能容纳的字符数（中文约 2 字符/英寸，英文约 4 字符/英寸）
    chars_per_line = box_width_in * 2.5
    estimated_lines = max(1, text_len / chars_per_line)

    # 估算能容纳的行数
    line_height_in = base_size / 72 * 1.4  # 行高 = 字号 * 行距
    max_lines = box_height_in / line_height_in

    if estimated_lines <= max_lines:
        # 内容能放下，可以用较大字号
        return min(base_size * 1.2, 24.0)  # 最大不超过 24pt
    else:
        # 内容较多，缩小字号以适应
        scale = max_lines / estimated_lines
        return max(base_size * scale, 12.0)  # 最小不低于 12pt


def _estimate_text_height(text: str, font_size: float, box_width_in: float) -> float:
    """估算文本所需的高度（英寸）"""
    chars_per_line = box_width_in * 2.5
    line_count = max(1, len(text) / chars_per_line)
    line_height = font_size / 72 * 1.4
    return line_count * line_height


def _adaptive_card_size(points: List[str], base_font: float, available_height: float) -> tuple:
    """根据要点数量和内容自适应计算卡片尺寸

    Returns:
        (card_height, font_size, cols, rows)
    """
    n = len(points)
    if n <= 3:
        cols, rows = n, 1
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 3, 3

    # 计算每张卡片可用高度
    card_h = min(available_height / rows - 0.3, 2.5)

    # 计算最大文本长度
    max_text_len = max(len(p) for p in points) if points else 10

    # 根据文本长度调整字号
    if max_text_len > 40:
        font_size = base_font * 0.85
    elif max_text_len > 25:
        font_size = base_font * 0.95
    else:
        font_size = base_font

    return card_h, font_size, cols, rows


class PPTSkillService:
    """基于 ppt-polished-deck-collab skill 脚本链路的 PPT 生成服务"""

    def __init__(self):
        self.output_dir = OUTPUT_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)

    async def create_ppt_from_report(
        self,
        report_content: str,
        template: str = "default",
        style: str = "professional",
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """
        主入口：从报告内容生成 PPT

        Args:
            report_content: Markdown 格式的报告内容
            template: 模板名称
            style: 视觉风格 (professional/creative/minimal)
            options: 额外选项 (canvas_format, custom_colors 等)

        Returns:
            生成的 PPTX 文件路径
        """
        options = options or {}
        start_time = time.time()

        # 提取标题
        title = self._extract_title(report_content)
        logger.info(f"[PPTSkillService] 开始生成 PPT | 标题: {title} | 风格: {style}")

        # Step 1: 初始化 workspace
        workspace_dir = self._init_workspace(title, options)
        logger.info(f"[PPTSkillService] Step 1: Workspace 初始化完成: {workspace_dir}")

        # Step 2: 写入源文件
        source_path = workspace_dir / "sources" / "report.md"
        source_path.parent.mkdir(parents=True, exist_ok=True)
        source_path.write_text(report_content, encoding="utf-8")

        # Step 3: LLM 生成 deck_narrative.md
        narrative_path = workspace_dir / "deck_narrative.md"
        await self._generate_narrative(workspace_dir, report_content, title, style, options)
        logger.info(f"[PPTSkillService] Step 3: Narrative 生成完成")

        # Step 4: 派生 slide_specs.yaml
        specs_path = self._derive_slide_specs(workspace_dir)
        logger.info(f"[PPTSkillService] Step 4: Slide specs 派生完成")

        # Step 5: 构建原生 PPTX
        pptx_path = self._build_pptx(workspace_dir, style, options)
        logger.info(f"[PPTSkillService] Step 5: PPTX 构建完成: {pptx_path}")

        # Step 6: 运行质量检查
        self._run_quality_gates(workspace_dir, pptx_path)
        logger.info(f"[PPTSkillService] Step 6: 质量检查完成")

        # Step 7: 复制到输出目录
        final_path = self._export_pptx(pptx_path, title, style)

        elapsed = time.time() - start_time
        logger.info(f"[PPTSkillService] PPT 生成完成 | 耗时: {elapsed:.1f}s | 路径: {final_path}")

        return str(final_path)

    def _extract_title(self, content: str) -> str:
        """从 Markdown 内容提取标题"""
        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("# "):
                return line[2:].strip()
        return "Untitled Deck"

    def _init_workspace(self, title: str, options: Dict[str, Any]) -> Path:
        """
        Step 1: 使用 init_deck_workspace.py 初始化 workspace

        调用 Skill 脚本创建标准目录结构和模板文件
        """
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[^\w一-鿿-]', '_', title)[:30]
        workspace_name = f"{safe_title}_{timestamp}"
        workspace_dir = self.output_dir / workspace_name

        cmd = [
            "python", str(SCRIPTS_DIR / "init_deck_workspace.py"),
            "--workspace-dir", str(workspace_dir),
            "--title", title,
        ]

        # 添加可选参数
        if options.get("audience"):
            cmd.extend(["--audience", options["audience"]])
        if options.get("scenario"):
            cmd.extend(["--scenario", options["scenario"]])
        if options.get("objective"):
            cmd.extend(["--objective", options["objective"]])

        try:
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=30,
                cwd=str(SKILL_DIR)
            )
            if result.returncode != 0:
                logger.warning(f"[PPTSkillService] init_deck_workspace 警告: {result.stderr}")
        except subprocess.TimeoutExpired:
            logger.warning("[PPTSkillService] init_deck_workspace 超时，手动创建目录")
            self._create_workspace_manually(workspace_dir, title)
        except FileNotFoundError:
            logger.warning("[PPTSkillService] init_deck_workspace 脚本不存在，手动创建目录")
            self._create_workspace_manually(workspace_dir, title)

        return workspace_dir

    def _create_workspace_manually(self, workspace_dir: Path, title: str):
        """手动创建 workspace 目录结构"""
        dirs = [
            "data", "assets/diagrams", "assets/charts", "assets/icons",
            "assets/images/prompts", "assets/images/generated", "assets/images/raw",
            "assets/tables", "build/generated", "build/pptx",
            "build/rendered/ppt_preview", "validation/package_preflight",
            "validation/structure_precheck", "validation/render_review",
            "validation/visual", "final", "sources"
        ]
        for d in dirs:
            (workspace_dir / d).mkdir(parents=True, exist_ok=True)

        # 写入最小 brief.md
        brief_path = workspace_dir / "brief.md"
        if not brief_path.exists():
            brief_path.write_text(
                f"# {title}\n\n## Deck Contract\n\n- audience: general\n- scenario: presentation\n- objective: understanding\n",
                encoding="utf-8"
            )

        # 写入最小 deck_narrative.md 模板
        narrative_path = workspace_dir / "deck_narrative.md"
        if not narrative_path.exists():
            narrative_path.write_text(
                f"---\ndeck:\n  title: \"{title}\"\n---\n\n# {title}\n\n",
                encoding="utf-8"
            )

    async def _generate_narrative(
        self,
        workspace_dir: Path,
        report_content: str,
        title: str,
        style: str,
        options: Dict[str, Any]
    ):
        """
        Step 3: 使用 LLM 生成 deck_narrative.md

        LLM 根据报告内容生成结构化的叙事文档，每页包含 slide_spec YAML 块
        """
        # 加载参考文档
        ref_files = self._load_reference_files()

        # 构建 prompt
        system_prompt = self._build_narrative_system_prompt(ref_files, title, style, options)
        user_prompt = self._build_narrative_user_prompt(report_content, title)

        # 调用 LLM
        client, model_name = _get_llm()

        for attempt in range(3):
            try:
                response = await client.chat.completions.create(
                    model=model_name,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    temperature=0.7,
                    max_tokens=12000
                )
                narrative_content = response.choices[0].message.content.strip()
                if narrative_content:
                    break
            except Exception as e:
                logger.warning(f"[PPTSkillService] LLM 调用失败 (尝试 {attempt + 1}/3): {e}")
                if attempt == 2:
                    raise RuntimeError(f"Narrative 生成失败: {e}")

        # 清理 markdown 代码块
        narrative_content = self._clean_markdown_fences(narrative_content)

        # 确保有 YAML frontmatter
        if not narrative_content.startswith("---"):
            narrative_content = self._add_frontmatter(narrative_content, title)

        # 写入 deck_narrative.md
        narrative_path = workspace_dir / "deck_narrative.md"
        narrative_path.write_text(narrative_content, encoding="utf-8")

        # 确保所有 YAML 块都有闭合的 ```
        self._fix_unclosed_yaml_blocks(narrative_path)

        logger.info(f"[PPTSkillService] deck_narrative.md 已写入: {narrative_path}")

    def _load_reference_files(self) -> Dict[str, str]:
        """加载 Skill 参考文档"""
        ref_files = {}
        ref_dir = REFERENCES_DIR

        # 核心参考文档
        core_files = {
            "schema_contract": "core/schema_contract.md",
            "style_profiles": "core/style_profiles.md",
            "principles": "core/principles.md",
        }

        # 工作流文档
        workflow_files = {
            "deck_workflow": "workflow/deck_workflow.md",
        }

        all_files = {**core_files, **workflow_files}

        for key, rel_path in all_files.items():
            file_path = ref_dir / rel_path
            if file_path.exists():
                try:
                    ref_files[key] = file_path.read_text(encoding="utf-8")
                except Exception as e:
                    logger.warning(f"[PPTSkillService] 无法读取 {file_path}: {e}")

        return ref_files

    def _build_narrative_system_prompt(
        self,
        ref_files: Dict[str, str],
        title: str,
        style: str,
        options: Dict[str, Any] = None
    ) -> str:
        """构建 narrative 生成的 system prompt"""

        options = options or {}
        schema_ref = ref_files.get("schema_contract", "")

        # 从 options 获取样式配置
        visual_profile = options.get("visual_profile", "corporate_clear")
        communication_profile = options.get("communication_profile", "business_report")
        density_profile = options.get("density_profile", "balanced_brief")
        delivery_context = options.get("delivery_context", "hybrid_review_deck")

        # 样式描述
        style_descriptions = {
            "visual_profile": {
                "corporate_clear": "商业清晰：清晰、稳定、克制的商务风格",
                "editorial_ink": "杂志墨韵：图文叙事、克制色彩、电子杂志风格",
                "swiss_modernist": "瑞士现代：网格秩序、单一强调色、发丝线、极简风格",
                "product_launch": "产品发布：hero页、产品图、路线图、发布会风格"
            },
            "communication_profile": {
                "business_report": "商业汇报：管理层同步、经营复盘、KPI图表",
                "technical_explainer": "技术说明：架构图、流程图、技术方案",
                "research_review": "研究评审：数据分析、图表展示、方法说明",
                "keynote_story": "演讲叙事：故事化、记忆点、视觉锚点"
            },
            "density_profile": {
                "dense_reference": "高密参考：完整标题、注释、来源、单位",
                "balanced_brief": "均衡适中：清晰结论与足够证据的平衡",
                "low_density_stage": "低密舞台：减少正文，信息压力转给讲者"
            },
            "delivery_context": {
                "self-contained_reading_deck": "自解释阅读：无人讲解也能看懂",
                "speaker-led_stage_deck": "舞台演讲：有人现场讲，页面低密度",
                "hybrid_review_deck": "混合审阅：讲完后转发，关键页自解释",
                "reference_or_appendix_deck": "参考附录：查阅型材料，高密度"
            }
        }

        visual_desc = style_descriptions["visual_profile"].get(visual_profile, visual_profile)
        comm_desc = style_descriptions["communication_profile"].get(communication_profile, communication_profile)
        density_desc = style_descriptions["density_profile"].get(density_profile, density_profile)
        delivery_desc = style_descriptions["delivery_context"].get(delivery_context, delivery_context)

        return f"""你是一个专业的 PPT 叙事设计师。你的任务是根据报告内容，生成一份结构化的 deck_narrative.md 文档。

## 样式配置（必须遵循）

- **视觉风格**：{visual_profile} - {visual_desc}
- **内容类型**：{communication_profile} - {comm_desc}
- **信息密度**：{density_profile} - {density_desc}
- **传播场景**：{delivery_context} - {delivery_desc}

## 输出格式要求

输出必须是一个完整的 Markdown 文档，包含：

1. **YAML frontmatter**（在 `---` 之间）：
```yaml
---
deck:
  title: "{title}"
  audience: "general"
  scenario: "presentation"
  objective: "understanding"
  source_context: "no_template"
  delivery_context: "{delivery_context}"
  communication_profile: "{communication_profile}"
  visual_profile: "{visual_profile}"
  density_profile: "{density_profile}"
  editability_profile: "fully_editable"
  theme_tokens:
    page_width_in: 13.333
    page_height_in: 7.5
    hero_title_font_pt: 48
    section_title_font_pt: 36
    page_title_font_pt: 28
    subtitle_font_pt: 20
    body_font_pt: 18
    label_font_pt: 14
    caption_font_pt: 12
    latin_font_name: "微软雅黑"
    east_asia_font_name: "微软雅黑"
    body_line_spacing_multiple: 1.4
    title_line_spacing_multiple: 1.2
---
```

2. **每页一个 `### Sxx | 标题` section**，包含：
   - 简短的自然语言描述（页面要传达什么）
   - 一个 ```yaml 代码块，**字段直接写在顶层，不要用 slide_spec: 包裹**，包含以下必需字段：

```yaml
title: "页面标题"
reader_question: "读者在这页想问的问题"
page_task: "这页要完成的任务"
reading_mode: "skim" 或 "read" 或 "study"
archetype: "cover" 或 "section_opener" 或 "key_insight" 或 "data_evidence" 或 "summary" 或 "comparison" 或 "timeline" 或 "process"
asset_mode: "text_only" 或 "chart_native" 或 "diagram" 或 "image" 或 "mixed"
validation_mode: "structure_only" 或 "full"
key_message: "这页的核心信息（一句话）"
layout_recipe: "full_bleed" 或 "two_column" 或 "three_column" 或 "grid" 或 "centered"
rhythm_role: "opener" 或 "evidence" 或 "breathing" 或 "transition" 或 "dense" 或 "closing"
```

**注意：YAML 块中不要写 `slide_spec:` 这个顶层 key，字段直接写！每个 YAML 块必须以 ``` 结束！**

3. **页面顺序**：
   - S01: 封面页 (archetype: cover)
   - S02: 目录/概览页 (archetype: section_opener)
   - S03-SN: 内容页（根据报告结构，每节一页或两页）
   - SN+1: 总结/结尾页 (archetype: summary)

## 参考规范

{schema_ref if schema_ref else "（无 schema_contract 参考）"}

## 重要规则

- 每页的 slide_spec YAML 块必须包含所有 8 个必需字段
- 页面数量控制在 6-20 页之间
- archetype 必须是上述枚举值之一
- 每页的 key_message 必须简洁有力，不超过 50 字
- 内容页的 key_message 应基于报告中的具体数据或结论

## narrative_markdown 内容要求（极其重要！）

在每个 YAML slide_spec 块**之后**，必须写该页的 **narrative_markdown**，这是页面的实际展示内容。要求：

1. **用 Markdown 无序列表（- ）写 3-6 个要点**，每个要点是一句话
2. **要点必须包含报告中的具体数据、数字、事实**，不要只写概括性描述
3. **格式示例**：

```yaml
title: "市场规模"
...
```

- 2026年中国AI产业规模预计突破8000亿元
- 大模型相关市场增速超过120%
- 企业级AI应用渗透率达到35%
- AI Agent市场规模约500亿元

**每个 slide_spec 块后面都必须跟这样的要点列表！包括封面页和总结页！这是 PPT 页面的实际内容来源。**
"""

    def _build_narrative_user_prompt(self, report_content: str, title: str) -> str:
        """构建 narrative 生成的 user prompt"""
        return f"""请根据以下报告内容，生成 deck_narrative.md 文档。

# 报告标题：{title}

# 报告内容：

{report_content}

---

请生成完整的 deck_narrative.md，包含 YAML frontmatter 和每页的 slide_spec YAML 块。
"""

    def _derive_slide_specs(self, workspace_dir: Path) -> Path:
        """
        Step 4: 使用 derive_slide_specs_from_narrative.py 派生 slide_specs.yaml

        从 deck_narrative.md 提取结构化的 slide specs
        """
        narrative_path = workspace_dir / "deck_narrative.md"
        specs_path = workspace_dir / "build" / "generated" / "slide_specs.yaml"
        specs_path.parent.mkdir(parents=True, exist_ok=True)

        cmd = [
            "python", str(SCRIPTS_DIR / "derive_slide_specs_from_narrative.py"),
            "--narrative", str(narrative_path),
            "--out-yaml", str(specs_path),
        ]

        try:
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=30,
                cwd=str(SKILL_DIR)
            )
            if result.returncode != 0:
                logger.warning(f"[PPTSkillService] derive_slide_specs 警告: {result.stderr}")
                # 如果脚本失败，尝试手动解析
                if not specs_path.exists():
                    self._derive_specs_manually(narrative_path, specs_path)
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            logger.warning(f"[PPTSkillService] derive_slide_specs 脚本执行失败: {e}")
            self._derive_specs_manually(narrative_path, specs_path)

        return specs_path

    def _derive_specs_manually(self, narrative_path: Path, specs_path: Path):
        """手动从 narrative 派生 slide specs（降级方案）"""
        try:
            import yaml
        except ImportError:
            logger.error("[PPTSkillService] PyYAML 未安装，无法派生 specs")
            raise RuntimeError("PyYAML 未安装")

        content = narrative_path.read_text(encoding="utf-8")

        # 解析 frontmatter
        frontmatter = {}
        fm_match = re.match(r"^---\n(.*?)\n---\n", content, re.DOTALL)
        if fm_match:
            try:
                frontmatter = yaml.safe_load(fm_match.group(1)) or {}
            except Exception:
                pass

        # 解析 slide sections
        slides = []
        slide_pattern = re.compile(r"^###\s+(S\d+)(?:\s*\|\s*(.+))?\s*$", re.MULTILINE)
        # 支持有/无闭合 ``` 的 YAML 块，兼容多种格式
        yaml_pattern = re.compile(r"```yaml(?:\s+slide_spec)?\s*\n(.*?)(?:\n\s*```|\Z)", re.DOTALL)

        matches = list(slide_pattern.finditer(content))
        for i, match in enumerate(matches):
            slide_id = match.group(1)
            title_hint = match.group(2).strip() if match.group(2) else f"Slide {slide_id}"

            # 提取 section 内容
            start = match.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(content)
            section = content[start:end].strip()

            # 提取 YAML block
            yaml_match = yaml_pattern.search(section)
            if yaml_match:
                try:
                    spec = yaml.safe_load(yaml_match.group(1))
                    if isinstance(spec, dict) and "slide_spec" in spec:
                        slide_data = spec["slide_spec"]
                        slide_data["slide_id"] = slide_id
                        slides.append(slide_data)
                        continue
                    elif isinstance(spec, dict) and "title" in spec:
                        # 直接是字段在顶层
                        spec["slide_id"] = slide_id
                        slides.append(spec)
                        continue
                except Exception:
                    pass

            # 降级：创建默认 spec
            slides.append({
                "slide_id": slide_id,
                "title": title_hint,
                "reader_question": "这页要传达什么？",
                "page_task": "展示关键信息",
                "reading_mode": "skim",
                "archetype": "key_insight",
                "asset_mode": "text_only",
                "validation_mode": "structure_only",
                "key_message": title_hint,
            })

        # 写入 YAML
        specs = {
            "deck": frontmatter.get("deck", {}),
            "slides": slides,
        }
        specs_path.write_text(
            yaml.dump(specs, allow_unicode=True, default_flow_style=False),
            encoding="utf-8"
        )

    def _build_pptx(
        self,
        workspace_dir: Path,
        style: str,
        options: Dict[str, Any]
    ) -> Path:
        """
        Step 5: 使用 python-pptx 构建原生 PPTX

        读取 slide_specs.yaml，使用 ppt_asset_helpers 构建可编辑的 PPTX
        """
        try:
            import yaml
        except ImportError:
            raise RuntimeError("PyYAML 未安装")

        # 导入 ppt_asset_helpers（使用绝对路径加载）
        try:
            import sys
            import importlib.util
            helpers_path = SCRIPTS_DIR / "ppt_asset_helpers.py"
            if not helpers_path.exists():
                raise RuntimeError(f"ppt_asset_helpers.py 不存在: {helpers_path}")
            spec = importlib.util.spec_from_file_location("ppt_asset_helpers", helpers_path)
            module = importlib.util.module_from_spec(spec)
            sys.modules["ppt_asset_helpers"] = module  # 注册到 sys.modules
            spec.loader.exec_module(module)
            new_presentation = module.new_presentation
            add_slide_header = module.add_slide_header
            add_text_block = module.add_text_block
            add_panel = module.add_panel
            default_palette = module.default_palette
            default_typography_tokens = module.default_typography_tokens
            save_presentation = module.save_presentation
        except Exception as e:
            logger.error(f"[PPTSkillService] 无法导入 ppt_asset_helpers: {e}")
            raise RuntimeError(f"无法导入 ppt_asset_helpers: {e}")

        # 读取 slide specs
        specs_path = workspace_dir / "build" / "generated" / "slide_specs.yaml"
        if not specs_path.exists():
            raise RuntimeError(f"slide_specs.yaml 不存在: {specs_path}")

        specs = yaml.safe_load(specs_path.read_text(encoding="utf-8"))
        slides = specs.get("slides", [])

        if not slides:
            raise RuntimeError("slide_specs.yaml 中没有 slides")

        # 从 narrative 提取每页的实际内容（narrative_markdown）
        narrative_path = workspace_dir / "deck_narrative.md"
        narrative_content_map = self._extract_narrative_content(narrative_path)
        for slide in slides:
            sid = slide.get("slide_id", "")
            if sid in narrative_content_map:
                slide["_narrative_content"] = narrative_content_map[sid]

        # 获取样式配置
        palette = default_palette()
        tokens = default_typography_tokens()

        # 覆盖为专业的中文 PPT 字号和字体
        tokens.update({
            "hero_title_font_pt": 48.0,
            "section_title_font_pt": 36.0,
            "page_title_font_pt": 28.0,
            "subtitle_font_pt": 20.0,
            "body_font_pt": 18.0,
            "label_font_pt": 14.0,
            "caption_font_pt": 12.0,
            "latin_font_name": "微软雅黑",
            "east_asia_font_name": "微软雅黑",
            "body_line_spacing_multiple": 1.4,
            "title_line_spacing_multiple": 1.2,
        })

        # 自定义颜色
        custom_colors = options.get("custom_colors") or {}
        if custom_colors.get("primary"):
            try:
                hex_color = custom_colors["primary"].lstrip("#")
                palette["title"] = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                palette["blue"] = palette["title"]
            except Exception:
                pass

        # 创建演示文稿
        prs = new_presentation()

        for i, slide_spec in enumerate(slides):
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            archetype = slide_spec.get("archetype", "key_insight")
            title = slide_spec.get("title", f"Slide {i+1}")
            key_message = slide_spec.get("key_message", "")
            figure_tag = f"S{i+1:02d}"

            if archetype == "cover":
                self._build_cover_slide(slide, slide_spec, palette, tokens, prs)
            elif archetype == "summary":
                self._build_summary_slide(slide, slide_spec, palette, tokens, prs)
            elif archetype == "section_opener":
                self._build_section_slide(slide, slide_spec, palette, tokens, prs)
            else:
                self._build_content_slide(slide, slide_spec, palette, tokens, prs, figure_tag)

        # 保存 PPTX
        pptx_dir = workspace_dir / "build" / "pptx"
        pptx_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = pptx_dir / "presentation.pptx"

        save_presentation(prs, pptx_path)

        return pptx_path

    def _build_cover_slide(self, slide, spec: Dict, palette: Dict, tokens: Dict, prs):
        """构建封面页"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return

        title = spec.get("title", "Untitled")
        key_message = spec.get("key_message", "")
        font_name = tokens.get("latin_font_name", "微软雅黑")

        # 背景色
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*palette["title"])

        # 主标题
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(13), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.alignment = PP_ALIGN.CENTER
        _set_font(title_para.font, font_name, tokens["hero_title_font_pt"], True, (255, 255, 255))

        # 副标题
        if key_message:
            sub_box = slide.shapes.add_textbox(Inches(2), Inches(5.0), Inches(12), Inches(1))
            sub_frame = sub_box.text_frame
            sub_frame.word_wrap = True
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = key_message
            sub_para.alignment = PP_ALIGN.CENTER
            _set_font(sub_para.font, font_name, tokens["subtitle_font_pt"], False, (200, 210, 225))

    def _build_section_slide(self, slide, spec: Dict, palette: Dict, tokens: Dict, prs):
        """构建章节页"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return

        title = spec.get("title", "")
        key_message = spec.get("key_message", "")
        font_name = tokens.get("latin_font_name", "微软雅黑")

        # 背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*palette["bg"])

        # 左侧色块
        accent = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0), Inches(0.4), Inches(9)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*palette["blue"])
        accent.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(11), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        _set_font(title_para.font, font_name, tokens["section_title_font_pt"], True, palette["title"])

        # 副标题
        if key_message:
            sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(11), Inches(0.8))
            sub_frame = sub_box.text_frame
            sub_frame.word_wrap = True
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = key_message
            _set_font(sub_para.font, font_name, tokens["subtitle_font_pt"], False, palette["subtitle"])

    def _build_summary_slide(self, slide, spec: Dict, palette: Dict, tokens: Dict, prs):
        """构建总结页 - 带 bullet points"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return

        title = spec.get("title", "Summary")
        key_message = spec.get("key_message", "")
        narrative_content = spec.get("_narrative_content", "")
        font_name = tokens.get("latin_font_name", "微软雅黑")

        # 解析要点
        points = self._parse_bullet_points(narrative_content)

        # 背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*palette["title"])

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(13), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.alignment = PP_ALIGN.CENTER
        _set_font(title_para.font, font_name, tokens["section_title_font_pt"], True, (255, 255, 255))

        # 核心信息
        if key_message:
            msg_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(0.8))
            msg_frame = msg_box.text_frame
            msg_frame.word_wrap = True
            msg_para = msg_frame.paragraphs[0]
            msg_para.text = key_message
            msg_para.alignment = PP_ALIGN.CENTER
            _set_font(msg_para.font, font_name, tokens["subtitle_font_pt"], False, (200, 210, 225))

        # 要点列表
        if points:
            content_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.8), Inches(11), Inches(4.5))
            tf = content_box.text_frame
            tf.word_wrap = True

            for idx, point in enumerate(points):
                if idx == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = f"  ●  {point}"
                _set_font(para.font, font_name, tokens["body_font_pt"] + 1, False, (220, 225, 235))
                para.line_spacing = tokens["body_line_spacing_multiple"]
                para.space_after = Pt(10)

    def _extract_narrative_content(self, narrative_path: Path) -> Dict[str, str]:
        """从 deck_narrative.md 提取每页的 narrative_markdown 内容"""
        if not narrative_path.exists():
            return {}

        content = narrative_path.read_text(encoding="utf-8")
        result = {}

        # 按 ### Sxx 分割
        slide_pattern = re.compile(r"^###\s+(S\d+)", re.MULTILINE)
        yaml_pattern = re.compile(r"```yaml(?:\s+slide_spec)?\s*\n(.*?)(?:\n\s*```|\Z)", re.DOTALL)

        matches = list(slide_pattern.finditer(content))
        for i, match in enumerate(matches):
            slide_id = match.group(1)
            start = match.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(content)
            section = content[start:end].strip()

            # 提取 YAML 块之后的内容作为 narrative_markdown
            yaml_match = yaml_pattern.search(section)
            if yaml_match:
                after_yaml = section[yaml_match.end():].strip()
                if after_yaml:
                    result[slide_id] = after_yaml

        return result

    def _parse_bullet_points(self, text: str) -> List[str]:
        """从 narrative_markdown 解析要点列表"""
        points = []
        for line in text.split("\n"):
            line = line.strip()
            # 匹配 - xxx 或 * xxx 或 1. xxx
            if re.match(r"^[-*]\s+", line):
                point = re.sub(r"^[-*]\s+", "", line).strip()
                if point:
                    points.append(point)
            elif re.match(r"^\d+[\.\)]\s+", line):
                point = re.sub(r"^\d+[\.\)]\s+", "", line).strip()
                if point:
                    points.append(point)
        return points

    def _build_content_slide(self, slide, spec: Dict, palette: Dict, tokens: Dict, prs, figure_tag: str):
        """构建内容页 - 根据 archetype 选择不同布局"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            return

        title = spec.get("title", "")
        key_message = spec.get("key_message", "")
        archetype = spec.get("archetype", "key_insight")
        narrative_content = spec.get("_narrative_content", "")
        font_name = tokens.get("latin_font_name", "微软雅黑")

        # 解析要点
        points = self._parse_bullet_points(narrative_content)
        if not points:
            # 降级：用 key_message 和 page_task
            points = [key_message, spec.get("page_task", "")]

        # 背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*palette["bg"])

        # 标题区
        title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.16), Inches(14.3), Inches(0.72))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = title
        _set_font(title_para.font, font_name, tokens["page_title_font_pt"], True, palette["title"])

        # 副标题（key_message）
        sub_box = slide.shapes.add_textbox(Inches(0.74), Inches(0.82), Inches(14.1), Inches(0.40))
        sub_para = sub_box.text_frame.paragraphs[0]
        sub_para.text = key_message
        _set_font(sub_para.font, font_name, tokens["subtitle_font_pt"], False, palette["subtitle"])

        # 图标标签
        tag_box = slide.shapes.add_textbox(Inches(14.85), Inches(8.43), Inches(0.46), Inches(0.18))
        tag_para = tag_box.text_frame.paragraphs[0]
        tag_para.text = figure_tag
        tag_para.alignment = PP_ALIGN.RIGHT
        _set_font(tag_para.font, font_name, tokens["caption_font_pt"], True, palette["muted"])

        # 根据 archetype 选择布局
        if archetype == "data_evidence":
            self._layout_data_evidence(slide, points, palette, tokens)
        elif archetype == "comparison":
            self._layout_comparison(slide, points, palette, tokens)
        elif archetype in ("timeline", "process"):
            self._layout_process(slide, points, palette, tokens)
        else:
            # 默认：要点列表布局
            self._layout_key_points(slide, points, palette, tokens)

    def _layout_key_points(self, slide, points: List[str], palette: Dict, tokens: Dict):
        """要点列表布局 - 带编号的要点卡片（自适应）"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            return

        font_name = tokens.get("latin_font_name", "微软雅黑")
        colors = [palette["blue"], palette["emerald"], palette["amber"],
                  palette["violet"], palette["rose"], palette["teal"]]

        # 自适应卡片尺寸
        available_h = 6.0  # 标题下方可用高度
        card_h, font_size, cols, rows = _adaptive_card_size(points, tokens["body_font_pt"], available_h)

        card_w = (14.0 - 0.25 * (cols - 1)) / cols
        start_x = 0.72
        start_y = 1.5

        for idx, point in enumerate(points[:9]):
            row = idx // cols
            col = idx % cols
            x = start_x + col * (card_w + 0.25)
            y = start_y + row * (card_h + 0.25)
            color = colors[idx % len(colors)]

            # 编号圆
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y + 0.05), Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(*color)
            circle.line.fill.background()
            _set_font(circle.text_frame.paragraphs[0].font, font_name, 14, True, (255, 255, 255))
            circle.text_frame.paragraphs[0].text = str(idx + 1)
            circle.text_frame.paragraphs[0].alignment = 1
            circle.text_frame.vertical_anchor = 1

            # 要点卡片
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.5), Inches(y), Inches(card_w - 0.5), Inches(card_h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*[min(255, c + 230) for c in [color[0]//10, color[1]//10, color[2]//10]])
            card.line.color.rgb = RGBColor(*[min(255, c + 180) for c in [color[0]//5, color[1]//5, color[2]//5]])
            card.line.width = Pt(0.5)

            text_frame = card.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.08)
            text_frame.margin_bottom = Inches(0.08)
            text_frame.vertical_anchor = 1

            para = text_frame.paragraphs[0]
            para.text = point
            # 自适应字号
            adaptive_size = _adaptive_font_size(point, font_size, card_w - 0.5, card_h)
            _set_font(para.font, font_name, adaptive_size, False, palette["title"])
            para.line_spacing = tokens["body_line_spacing_multiple"]

    def _layout_data_evidence(self, slide, points: List[str], palette: Dict, tokens: Dict):
        """数据证据布局 - 自适应"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            return

        font_name = tokens.get("latin_font_name", "微软雅黑")
        colors = [palette["blue"], palette["emerald"], palette["amber"],
                  palette["violet"], palette["rose"], palette["teal"]]

        # 提取含数字的要点作为重点数据
        data_points = [p for p in points if any(c.isdigit() for c in p)]
        other_points = [p for p in points if not any(c.isdigit() for c in p)]

        # 自适应计算数据卡片高度
        n_data = min(len(data_points), 3)
        if n_data > 0:
            data_card_h = min(5.5 / n_data - 0.2, 1.8)
        else:
            data_card_h = 1.6

        # 右侧：数据高亮卡片
        for idx, dp in enumerate(data_points[:3]):
            y = 1.5 + idx * (data_card_h + 0.2)
            color = colors[idx % len(colors)]
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(9.5), Inches(y), Inches(5.5), Inches(data_card_h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*color)
            card.line.fill.background()

            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)
            tf.vertical_anchor = 1

            para = tf.paragraphs[0]
            para.text = dp
            # 自适应字号
            adaptive_size = _adaptive_font_size(dp, tokens["body_font_pt"] + 2, 5.2, data_card_h)
            _set_font(para.font, font_name, adaptive_size, True, (255, 255, 255))
            para.line_spacing = 1.2

        # 左侧：非数据要点列表
        if other_points:
            left_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.5), Inches(8.3), Inches(6.0))
            text_frame = left_box.text_frame
            text_frame.word_wrap = True

            # 自适应字号
            max_len = max(len(p) for p in other_points) if other_points else 20
            adaptive_size = _adaptive_font_size("x" * max_len, tokens["body_font_pt"], 8.0, 5.5)

            for idx, point in enumerate(other_points):
                if idx == 0:
                    para = text_frame.paragraphs[0]
                else:
                    para = text_frame.add_paragraph()
                para.text = f"● {point}"
                _set_font(para.font, font_name, adaptive_size, False, palette["title"])
                para.line_spacing = tokens["body_line_spacing_multiple"]
                para.space_after = Pt(8)

    def _layout_comparison(self, slide, points: List[str], palette: Dict, tokens: Dict):
        """对比布局 - 左右两栏对比（自适应）"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            return

        font_name = tokens.get("latin_font_name", "微软雅黑")
        mid = len(points) // 2
        left_points = points[:mid] if mid > 0 else points[:1]
        right_points = points[mid:] if mid > 0 else points[1:]

        # 自适应字号
        all_points = left_points + right_points
        max_len = max(len(p) for p in all_points) if all_points else 20
        adaptive_size = _adaptive_font_size("x" * max_len, tokens["body_font_pt"], 6.5, 5.5)

        # 左栏
        left_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.72), Inches(1.5), Inches(7.0), Inches(6.0)
        )
        left_panel.fill.solid()
        left_panel.fill.fore_color.rgb = RGBColor(240, 245, 255)
        left_panel.line.color.rgb = RGBColor(*palette["blue"])
        left_panel.line.width = Pt(1)

        left_tf = left_panel.text_frame
        left_tf.word_wrap = True
        left_tf.margin_left = Inches(0.2)
        left_tf.margin_top = Inches(0.15)

        for idx, point in enumerate(left_points):
            if idx == 0:
                para = left_tf.paragraphs[0]
            else:
                para = left_tf.add_paragraph()
            para.text = f"● {point}"
            _set_font(para.font, font_name, adaptive_size, False, palette["title"])
            para.line_spacing = tokens["body_line_spacing_multiple"]
            para.space_after = Pt(6)

        # 右栏
        right_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.2), Inches(1.5), Inches(7.0), Inches(6.0)
        )
        right_panel.fill.solid()
        right_panel.fill.fore_color.rgb = RGBColor(240, 255, 245)
        right_panel.line.color.rgb = RGBColor(*palette["emerald"])
        right_panel.line.width = Pt(1)

        right_tf = right_panel.text_frame
        right_tf.word_wrap = True
        right_tf.margin_left = Inches(0.2)
        right_tf.margin_top = Inches(0.15)

        for idx, point in enumerate(right_points):
            if idx == 0:
                para = right_tf.paragraphs[0]
            else:
                para = right_tf.add_paragraph()
            para.text = f"● {point}"
            _set_font(para.font, font_name, adaptive_size, False, palette["title"])
            para.line_spacing = tokens["body_line_spacing_multiple"]
            para.space_after = Pt(6)

    def _layout_process(self, slide, points: List[str], palette: Dict, tokens: Dict):
        """流程/时间线布局 - 横向步骤（自适应）"""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            return

        font_name = tokens.get("latin_font_name", "微软雅黑")
        n = min(len(points), 5)
        step_w = (14.5 - 0.3 * (n - 1)) / n
        start_x = 0.72

        # 自适应字号
        max_len = max(len(p) for p in points[:n]) if points else 20
        adaptive_size = _adaptive_font_size("x" * max_len, tokens["body_font_pt"], step_w - 0.3, 3.5)

        for idx in range(n):
            x = start_x + idx * (step_w + 0.3)
            color = [palette["blue"], palette["emerald"], palette["amber"],
                     palette["violet"], palette["rose"]][idx % 5]

            # 步骤卡片
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2.0), Inches(step_w), Inches(5.0)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*[min(255, c + 235) for c in [color[0]//10, color[1]//10, color[2]//10]])
            card.line.color.rgb = RGBColor(*color)
            card.line.width = Pt(1.5)

            # 步骤编号
            num_box = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + step_w/2 - 0.25), Inches(2.2), Inches(0.5), Inches(0.5)
            )
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = RGBColor(*color)
            num_box.line.fill.background()
            _set_font(num_box.text_frame.paragraphs[0].font, font_name, 16, True, (255, 255, 255))
            num_box.text_frame.paragraphs[0].text = str(idx + 1)
            num_box.text_frame.paragraphs[0].alignment = 1
            num_box.text_frame.vertical_anchor = 1

            # 步骤内容
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.9), Inches(step_w - 0.2), Inches(3.8)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = points[idx]
            _set_font(para.font, font_name, adaptive_size, False, palette["title"])
            para.line_spacing = 1.3
            para.alignment = 1

    def _run_quality_gates(self, workspace_dir: Path, pptx_path: Path):
        """
        Step 6: 运行质量检查

        调用 Skill 的 package_preflight 和 structure_precheck 脚本
        """
        # Gate 1: Package Preflight
        try:
            cmd1 = [
                "python", str(SCRIPTS_DIR / "check_pptx_package_preflight.py"),
                "--pptx", str(pptx_path),
                "--workspace-dir", str(workspace_dir),
                "--fail-on", "warning",
            ]
            result1 = subprocess.run(
                cmd1, capture_output=True, text=True, timeout=60,
                cwd=str(SKILL_DIR)
            )
            if result1.returncode != 0:
                logger.warning(f"[PPTSkillService] package_preflight 警告:\n{result1.stdout}")
            else:
                logger.info("[PPTSkillService] package_preflight 通过")
        except Exception as e:
            logger.warning(f"[PPTSkillService] package_preflight 执行失败: {e}")

        # Gate 2: Structure Precheck
        try:
            cmd2 = [
                "python", str(SCRIPTS_DIR / "check_pptx_structure_precheck.py"),
                "--pptx", str(pptx_path),
                "--workspace-dir", str(workspace_dir),
                "--fail-on", "warning",
            ]
            result2 = subprocess.run(
                cmd2, capture_output=True, text=True, timeout=60,
                cwd=str(SKILL_DIR)
            )
            if result2.returncode != 0:
                logger.warning(f"[PPTSkillService] structure_precheck 警告:\n{result2.stdout}")
            else:
                logger.info("[PPTSkillService] structure_precheck 通过")
        except Exception as e:
            logger.warning(f"[PPTSkillService] structure_precheck 执行失败: {e}")

    def _export_pptx(self, pptx_path: Path, title: str, style: str) -> Path:
        """Step 7: 复制 PPTX 到最终输出目录"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[^\w一-鿿-]', '_', title)[:20]
        final_name = f"pptskill_{safe_title}_{style}_{timestamp}.pptx"
        final_path = self.output_dir / final_name

        shutil.copy2(str(pptx_path), str(final_path))
        logger.info(f"[PPTSkillService] PPTX 已复制到: {final_path}")

        return final_path

    def _fix_unclosed_yaml_blocks(self, narrative_path: Path):
        """修复未闭合的 YAML 代码块"""
        content = narrative_path.read_text(encoding="utf-8")

        # 找到所有 ```yaml 开头的块，如果没有对应的 ```，则在块末尾添加
        lines = content.split('\n')
        result = []
        in_yaml_block = False

        for line in lines:
            if line.strip().startswith('```yaml'):
                in_yaml_block = True
                result.append(line)
            elif in_yaml_block and line.strip() == '```':
                in_yaml_block = False
                result.append(line)
            elif in_yaml_block and line.strip().startswith('### '):
                # 遇到下一个 section，说明上一个 YAML 块没有闭合
                result.append('```')
                result.append('')
                in_yaml_block = False
                result.append(line)
            else:
                result.append(line)

        # 如果文件结束时 YAML 块仍未闭合
        if in_yaml_block:
            result.append('```')

        narrative_path.write_text('\n'.join(result), encoding="utf-8")
        logger.info("[PPTSkillService] 已修复未闭合的 YAML 块")

    def _clean_markdown_fences(self, content: str) -> str:
        """清理 markdown 代码块标记"""
        # 移除开头的 ```markdown 或 ```
        content = re.sub(r'^```(?:markdown)?\s*\n', '', content)
        # 移除结尾的 ```
        content = re.sub(r'\n```\s*$', '', content)
        return content.strip()

    def _add_frontmatter(self, content: str, title: str) -> str:
        """为缺少 frontmatter 的内容添加默认 frontmatter"""
        frontmatter = f"""---
deck:
  title: "{title}"
  audience: "general"
  scenario: "presentation"
  objective: "understanding"
  source_context: "no_template"
  delivery_context: "hybrid_review_deck"
  communication_profile: "business_report"
  visual_profile: "corporate_clear"
  density_profile: "balanced_brief"
  editability_profile: "fully_editable"
  theme_tokens:
    page_width_in: 13.333
    page_height_in: 7.5
    hero_title_font_pt: 48
    section_title_font_pt: 36
    page_title_font_pt: 28
    subtitle_font_pt: 20
    body_font_pt: 18
    label_font_pt: 14
    caption_font_pt: 12
    latin_font_name: "微软雅黑"
    east_asia_font_name: "微软雅黑"
    body_line_spacing_multiple: 1.4
    title_line_spacing_multiple: 1.2
---

"""
        return frontmatter + content
