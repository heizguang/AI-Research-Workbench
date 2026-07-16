"""
生成专业PPT模板
运行一次即可生成模板文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os


def create_professional_template(output_path: str = "./templates/professional.pptx"):
    """创建专业商务模板"""
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs = Presentation()

    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 主色调
    PRIMARY = RGBColor(0, 82, 155)      # 深蓝
    SECONDARY = RGBColor(0, 120, 215)   # 亮蓝
    ACCENT = RGBColor(255, 152, 0)      # 橙色
    WHITE = RGBColor(255, 255, 255)
    DARK = RGBColor(33, 33, 33)
    GRAY = RGBColor(117, 117, 117)

    # ===== 封面页 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白

    # 底部色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(0), prs.slide_width, Inches(3.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    # 装饰线
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1.5), Inches(3.35), Inches(2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # 主标题
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "{{TITLE}}"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "{{SUBTITLE}}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(200, 220, 240)
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.LEFT

    # 底部信息
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "多智能体报告生成系统"
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY
    p3.font.name = "Microsoft YaHei"
    p3.alignment = PP_ALIGN.LEFT

    # ===== 目录页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # 顶部色条
    shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    # 左侧色条
    shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0.08), Inches(0.4), Inches(7.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY
    shape.line.fill.background()

    # 目录标题
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{Toc}}"
    p.font.size = Pt(28)
    p.font.color.rgb = PRIMARY
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    # 目录内容
    txBox2 = slide2.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "{{TocContent}}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK
    p2.font.name = "Microsoft YaHei"

    # ===== 内容页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # 顶部色条
    shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    # 左侧色条
    shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0.08), Inches(0.4), Inches(7.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY
    shape.line.fill.background()

    # 标题
    txBox = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{ContentTitle}}"
    p.font.size = Pt(28)
    p.font.color.rgb = PRIMARY
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    # 标题下划线
    shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(1), Inches(1.45), Inches(1.5), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # 内容
    txBox2 = slide3.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "{{Content}}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK
    p2.font.name = "Microsoft YaHei"

    # 页码
    txBox3 = slide3.shapes.add_textbox(Inches(12), Inches(6.8), Inches(1), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "{{Page}}"
    p3.font.size = Pt(11)
    p3.font.color.rgb = GRAY
    p3.font.name = "Microsoft YaHei"
    p3.alignment = PP_ALIGN.RIGHT

    # ===== 结束页 =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # 全屏背景
    shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    # 装饰线
    shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(5.5), Inches(2.8), Inches(2.3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # 感谢标题
    txBox = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{EndTitle}}"
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    # 副内容
    txBox2 = slide4.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.333), Inches(1.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "{{EndContent}}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 220, 240)
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER

    # 底部
    txBox3 = slide4.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "多智能体报告生成系统"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(150, 170, 190)
    p3.font.name = "Microsoft YaHei"
    p3.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    print(f"✅ 模板已生成: {output_path}")


if __name__ == "__main__":
    create_professional_template()
