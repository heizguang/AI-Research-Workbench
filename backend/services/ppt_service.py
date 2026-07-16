"""
PPT服务
处理PPT的生成和导出 - 专业汇报级PPT
"""

import os
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class PPTService:
    """PPT服务类 - 生成专业汇报级PPT"""

    def __init__(self, output_dir: str = "./data/ppt"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    async def create_ppt(
        self,
        ppt_content: Dict[str, Any],
        template: str = "default",
        style: str = "professional"
    ) -> str:
        """创建专业PPT文件"""
        prs = Presentation()

        # 16:9 宽屏
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        style_config = self._get_style_config(style)
        slides_data = ppt_content.get("slides", [])

        # 逐页添加幻灯片
        for i, slide_data in enumerate(slides_data):
            is_first = (i == 0)
            is_last = (i == len(slides_data) - 1)

            if is_first:
                self._add_cover_slide(prs, slide_data, style_config)
            elif is_last:
                self._add_ending_slide(prs, slide_data, style_config)
            else:
                self._add_content_slide(prs, slide_data, style_config, i, len(slides_data))

        # 保存文件
        filename = f"presentation_{template}_{style}.pptx"
        file_path = os.path.join(self.output_dir, filename)
        prs.save(file_path)

        return file_path

    def _get_style_config(self, style: str) -> Dict[str, Any]:
        """获取样式配置"""
        styles = {
            "professional": {
                "primary_color": RGBColor(0, 82, 155),      # 主色：深蓝
                "secondary_color": RGBColor(0, 120, 215),    # 副色：亮蓝
                "accent_color": RGBColor(255, 152, 0),       # 强调色：橙色
                "dark_text": RGBColor(33, 33, 33),           # 深色文字
                "light_text": RGBColor(255, 255, 255),       # 浅色文字
                "gray_text": RGBColor(117, 117, 117),        # 灰色文字
                "bg_color": RGBColor(255, 255, 255),         # 背景色
                "light_bg": RGBColor(245, 247, 250),         # 浅背景色
                "title_size": Pt(36),
                "subtitle_size": Pt(20),
                "heading_size": Pt(28),
                "body_size": Pt(16),
                "caption_size": Pt(12),
            },
            "creative": {
                "primary_color": RGBColor(255, 87, 34),
                "secondary_color": RGBColor(255, 152, 0),
                "accent_color": RGBColor(76, 175, 80),
                "dark_text": RGBColor(33, 33, 33),
                "light_text": RGBColor(255, 255, 255),
                "gray_text": RGBColor(117, 117, 117),
                "bg_color": RGBColor(255, 255, 255),
                "light_bg": RGBColor(255, 249, 245),
                "title_size": Pt(36),
                "subtitle_size": Pt(20),
                "heading_size": Pt(28),
                "body_size": Pt(16),
                "caption_size": Pt(12),
            },
            "minimal": {
                "primary_color": RGBColor(33, 33, 33),
                "secondary_color": RGBColor(66, 66, 66),
                "accent_color": RGBColor(200, 200, 200),
                "dark_text": RGBColor(33, 33, 33),
                "light_text": RGBColor(255, 255, 255),
                "gray_text": RGBColor(150, 150, 150),
                "bg_color": RGBColor(255, 255, 255),
                "light_bg": RGBColor(248, 248, 248),
                "title_size": Pt(34),
                "subtitle_size": Pt(18),
                "heading_size": Pt(26),
                "body_size": Pt(15),
                "caption_size": Pt(12),
            }
        }
        return styles.get(style, styles["professional"])

    def _add_cover_slide(self, prs: Presentation, slide_data: Dict, config: Dict):
        """添加封面页 - 全色背景 + 大标题"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 底部色块
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), prs.slide_width, Inches(3.2),
                        config["primary_color"])

        # 装饰线
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(1.5), Inches(3.35), Inches(2), Pt(4),
                        config["accent_color"])

        # 主标题（居中偏上）
        title_text = slide_data.get("title", "汇报PPT")
        content = slide_data.get("content", [])

        textbox = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(1.5))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.color.rgb = config["light_text"]
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.LEFT

        # 副标题（如果有）
        if content:
            textbox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10), Inches(0.8))
            tf2 = textbox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = content[0]
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(200, 220, 240)
            p2.font.name = "Microsoft YaHei"
            p2.alignment = PP_ALIGN.LEFT

        # 底部日期信息
        textbox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
        tf3 = textbox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "多智能体报告生成系统 · AI Generated"
        p3.font.size = Pt(14)
        p3.font.color.rgb = config["gray_text"]
        p3.font.name = "Microsoft YaHei"
        p3.alignment = PP_ALIGN.LEFT

        # 添加备注
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _add_content_slide(self, prs: Presentation, slide_data: Dict, config: Dict,
                           index: int, total: int):
        """添加内容页 - 左侧色条 + 标题 + 要点"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        title = slide_data.get("title", "")
        content = slide_data.get("content", [])
        notes = slide_data.get("notes", "")

        # 顶部色条
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), prs.slide_width, Inches(0.08),
                        config["primary_color"])

        # 左侧装饰色块
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0.08), Inches(0.4), Inches(7.42),
                        config["secondary_color"])

        # 页码
        page_textbox = slide.shapes.add_textbox(Inches(12), Inches(6.8), Inches(1), Inches(0.5))
        page_tf = page_textbox.text_frame
        page_p = page_tf.paragraphs[0]
        page_p.text = f"{index}/{total - 2}"
        page_p.font.size = Pt(11)
        page_p.font.color.rgb = config["gray_text"]
        page_p.font.name = "Microsoft YaHei"
        page_p.alignment = PP_ALIGN.RIGHT

        # 标题区域（带下划线装饰）
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(0.9))
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = config["heading_size"]
        title_p.font.color.rgb = config["primary_color"]
        title_p.font.bold = True
        title_p.font.name = "Microsoft YaHei"

        # 标题下划线
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(1), Inches(1.45), Inches(1.5), Pt(3),
                        config["accent_color"])

        # 内容区域 - 要点列表
        content_top = Inches(1.8)
        content_left = Inches(1.2)
        content_width = Inches(10.5)

        if content:
            textbox = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(5))
            tf = textbox.text_frame
            tf.word_wrap = True

            for j, point in enumerate(content):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                # 添加项目符号
                p.text = f"  •  {point}"
                p.font.size = config["body_size"]
                p.font.color.rgb = config["dark_text"]
                p.font.name = "Microsoft YaHei"
                p.space_after = Pt(12)
                p.space_before = Pt(4)

                # 设置行间距
                p.line_spacing = Pt(28)

        # 添加备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_ending_slide(self, prs: Presentation, slide_data: Dict, config: Dict):
        """添加结束页 - 感谢页"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        title = slide_data.get("title", "感谢")
        content = slide_data.get("content", [])

        # 全屏背景色
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), prs.slide_width, prs.slide_height,
                        config["primary_color"])

        # 装饰线
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(5.5), Inches(2.8), Inches(2.3), Pt(3),
                        config["accent_color"])

        # 主标题
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.2))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.color.rgb = config["light_text"]
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER

        # 副内容
        if content:
            textbox2 = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.333), Inches(1.5))
            tf2 = textbox2.text_frame
            for i, line in enumerate(content):
                if i == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                p2.text = line
                p2.font.size = Pt(18)
                p2.font.color.rgb = RGBColor(200, 220, 240)
                p2.font.name = "Microsoft YaHei"
                p2.alignment = PP_ALIGN.CENTER

        # 底部
        textbox3 = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.5))
        tf3 = textbox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "多智能体报告生成系统 · AI Generated"
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(150, 170, 190)
        p3.font.name = "Microsoft YaHei"
        p3.alignment = PP_ALIGN.CENTER

        # 添加备注
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _add_shape(self, slide, shape_type, left, top, width, height, color):
        """添加装饰形状"""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()  # 无边框
        return shape

    async def get_ppt_template(self, template_name: str) -> Dict[str, Any]:
        """获取PPT模板信息"""
        templates = {
            "default": {
                "name": "默认模板",
                "description": "简洁专业的默认模板",
                "style": "professional"
            },
            "business": {
                "name": "商务模板",
                "description": "适合商务汇报的模板",
                "style": "professional"
            },
            "creative": {
                "name": "创意模板",
                "description": "富有创意的设计模板",
                "style": "creative"
            },
            "minimal": {
                "name": "简约模板",
                "description": "简约风格的模板",
                "style": "minimal"
            }
        }
        return templates.get(template_name, templates["default"])

    async def list_templates(self) -> List[Dict[str, Any]]:
        """列出所有可用模板"""
        return [
            await self.get_ppt_template("default"),
            await self.get_ppt_template("business"),
            await self.get_ppt_template("creative"),
            await self.get_ppt_template("minimal")
        ]
