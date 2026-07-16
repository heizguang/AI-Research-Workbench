# -*- coding: utf-8 -*-
"""
本地测试 PPTSkillService - 不提交到仓库
"""
import asyncio
import sys
from pathlib import Path

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT / "backend"))

from dotenv import load_dotenv
load_dotenv(ROOT / "backend" / ".env")


TEST_REPORT = """# 2026年中国AI产业发展报告

## 1. 执行摘要
2026年上半年，中国AI产业进入规模化落地阶段。大模型从技术竞争转向应用竞争，Agent和多模态成为主要方向。国产大模型在中文理解、代码生成等核心能力上已接近国际先进水平。

## 2. 市场规模
- 2026年中国AI产业规模预计突破8000亿元
- 大模型相关市场增速超过120%
- 企业级AI应用渗透率达到35%
- AI Agent市场规模约500亿元

## 3. 技术趋势
- 多模态大模型成为标配，文本、图像、视频、音频统一理解
- Agent框架成熟度提升，从单步执行到多步规划
- 端侧部署加速，手机、PC、IoT设备广泛集成AI能力
- RAG架构向长上下文演进，200万Token窗口开始替代传统检索

## 4. 竞争格局
- 第一梯队：百度文心、阿里通义、字节豆包、腾讯混元
- 第二梯队：智谱、月之暗面、MiniMax、零一万物
- 国际竞争：GPT-5.5/Claude/Gemini同期快速迭代
- 差异化方向：垂直行业、端侧部署、成本优化

## 5. 应用场景
- 智能客服：覆盖率从40%提升至70%
- 代码辅助：开发者使用率超过60%
- 内容创作：图文视频生成进入商业化阶段
- 金融分析：研报自动生成、风险预警

## 6. 挑战与风险
- 算力成本仍是最大瓶颈
- 数据合规要求趋严
- AI幻觉问题尚未根本解决
- 人才竞争加剧

## 7. 展望
2026年下半年，AI产业将从"能用"走向"好用"，Agent和垂直场景将成为主要增长点。预计2027年AI将深度融入企业核心业务流程。
"""


async def main():
    from services.ppt.ppt_skill_service import PPTSkillService

    service = PPTSkillService()

    print("=" * 60)
    print("  PPTSkillService 本地测试")
    print("=" * 60)
    print(f"\n报告长度: {len(TEST_REPORT)} 字符")
    print(f"报告标题: 2026年中国AI产业发展报告\n")

    try:
        pptx_path = await service.create_ppt_from_report(
            report_content=TEST_REPORT,
            template="default",
            style="professional",
            options={
                "canvas_format": "ppt169",
                "audience": "企业决策者、技术管理者",
            }
        )

        print(f"\n{'=' * 60}")
        print(f"  生成成功!")
        print(f"{'=' * 60}")
        print(f"  PPTX 路径: {pptx_path}")

        pptx_file = Path(pptx_path)
        if pptx_file.exists():
            size_kb = pptx_file.stat().st_size / 1024
            print(f"  文件大小: {size_kb:.1f} KB")

            # 验证 PPTX 内容
            try:
                from pptx import Presentation
                prs = Presentation(str(pptx_file))
                print(f"  幻灯片数: {len(prs.slides)}")
                for i, slide in enumerate(prs.slides):
                    shapes = len(slide.shapes)
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text_frame.text.strip()[:40]
                            if text:
                                texts.append(text)
                    preview = texts[0] if texts else "(无文本)"
                    print(f"    S{i+1:02d}: {shapes} 个形状 | {preview}")
            except Exception as e:
                print(f"  验证失败: {e}")
        else:
            print(f"  [错误] 文件不存在!")

    except Exception as e:
        print(f"\n{'=' * 60}")
        print(f"  生成失败!")
        print(f"{'=' * 60}")
        print(f"  错误: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    asyncio.run(main())
